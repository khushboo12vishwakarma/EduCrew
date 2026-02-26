from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from PIL import Image
import requests
import logging


logger = logging.getLogger(__name__)


class PowerPointGenerator:
    """
    Professional College-Ready Presentation Generator with Dynamic Image Support
    
    Generates presentations with:
    - Professional styling and colors
    - Title slide with presenter info
    - Agenda slide
    - Content slides with images
    - Conclusion and thank you slides
    """
    
    def __init__(self, output_path, theme="college_blue"):
        """
        Initialize PowerPoint Generator
        
        Args:
            output_path: Path where PPTX file will be saved
            theme: Color theme ('college_blue' or 'college_teal')
        """
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # PROFESSIONAL COLOR THEMES
        self.themes = {
            "college_blue": {
                "primary_dark": RGBColor(13, 45, 89),      # Navy Blue
                "primary_light": RGBColor(230, 240, 250),  # Light Blue
                "accent": RGBColor(0, 120, 215),           # Bright Blue
                "accent_light": RGBColor(200, 230, 255),   # Light Accent
                "white": RGBColor(255, 255, 255),
                "dark_text": RGBColor(30, 30, 30),
                "medium_text": RGBColor(70, 70, 70),
                "light_text": RGBColor(255, 255, 255),
            },
            "college_teal": {
                "primary_dark": RGBColor(0, 102, 102),
                "primary_light": RGBColor(230, 245, 245),
                "accent": RGBColor(0, 150, 150),
                "accent_light": RGBColor(200, 240, 240),
                "white": RGBColor(255, 255, 255),
                "dark_text": RGBColor(30, 30, 30),
                "medium_text": RGBColor(70, 70, 70),
                "light_text": RGBColor(255, 255, 255),
            }
        }
        
        self.theme = self.themes.get(theme, self.themes["college_blue"])
        self.font_name = "Calibri"
        
        logger.info(f"[PPT-INIT] PowerPointGenerator initialized with theme: {theme}")

    def _download_image(self, url):
        """
        Download image from URL and convert to BytesIO
        
        This method:
        - Downloads image from Pixabay/Unsplash URL
        - Validates it's a real image
        - Returns BytesIO object for embedding
        
        Args:
            url: Image URL from Pixabay or other source
        
        Returns:
            BytesIO object with image data, or None if failed
        """
        if not url:
            logger.debug("[IMG-DL] Empty URL provided")
            return None
        
        try:
            url_preview = str(url)[:70]
            logger.debug(f"[IMG-DL] Downloading: {url_preview}...")
            
            # Download with timeout
            response = requests.get(url, timeout=15, stream=True)
            response.raise_for_status()
            
            logger.debug(f"[IMG-DL] Response status: {response.status_code}")
            
            # Read into BytesIO
            img_bytes = BytesIO(response.content)
            
            logger.debug(f"[IMG-DL] Downloaded {len(response.content)} bytes")
            
            # Validate it's a real image
            img = Image.open(img_bytes)
            img.verify()
            
            # Re-open after verify
            img_bytes.seek(0)
            
            logger.info(f"[IMG-DL] ✅ Image downloaded: {len(response.content)} bytes")
            return img_bytes
            
        except requests.exceptions.Timeout:
            logger.warning(f"[IMG-DL] ⏱️  Timeout downloading image")
            return None
        except requests.exceptions.ConnectionError as e:
            logger.warning(f"[IMG-DL] 🔌 Connection error: {e}")
            return None
        except Exception as e:
            logger.warning(f"[IMG-DL] ❌ Error: {str(e)}")
            return None

    def _add_background(self, slide, color):
        """
        Add solid background color to slide
        
        Args:
            slide: Slide object
            color: RGBColor object
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _title_slide(self, title, presenter_name, presentation_date, subject):
        """
        Create professional title slide
        
        Layout:
        - Dark background
        - Large title with accent color
        - Subject line
        - Presenter name
        - Date
        
        Args:
            title: Presentation title
            presenter_name: Name of presenter
            presentation_date: Date of presentation
            subject: Subject or course name
        """
        logger.info(f"[PPT-TITLE] Creating title slide: '{title}'")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide, self.theme["primary_dark"])

        # Top accent bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.theme["accent"]
        top_bar.line.fill.background()

        # Main title - LARGE, CENTERED, ACCENT COLOR
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(52)
        title_p.font.bold = True
        title_p.font.color.rgb = self.theme["accent"]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.name = self.font_name

        # Divider line
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2), Inches(2.8), Inches(6), Inches(0.03)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.theme["accent"]
        divider.line.fill.background()

        # Subject line
        if subject:
            subject_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
            subject_frame = subject_box.text_frame
            subject_p = subject_frame.paragraphs[0]
            subject_p.text = subject
            subject_p.font.size = Pt(16)
            subject_p.font.italic = True
            subject_p.font.color.rgb = self.theme["accent_light"]
            subject_p.alignment = PP_ALIGN.CENTER
            subject_p.font.name = self.font_name

        # "Presented by" label
        presented_label_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.25))
        presented_label_frame = presented_label_box.text_frame
        presented_label_p = presented_label_frame.paragraphs[0]
        presented_label_p.text = "Presented by"
        presented_label_p.font.size = Pt(12)
        presented_label_p.font.color.rgb = self.theme["light_text"]
        presented_label_p.alignment = PP_ALIGN.CENTER
        presented_label_p.font.name = self.font_name

        # Presenter name
        presenter_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.55), Inches(9), Inches(0.5))
        presenter_frame = presenter_box.text_frame
        presenter_p = presenter_frame.paragraphs[0]
        presenter_p.text = presenter_name
        presenter_p.font.size = Pt(22)
        presenter_p.font.bold = True
        presenter_p.font.color.rgb = self.theme["white"]
        presenter_p.alignment = PP_ALIGN.CENTER
        presenter_p.font.name = self.font_name

        # Date at bottom
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
        date_frame = date_box.text_frame
        date_p = date_frame.paragraphs[0]
        date_p.text = presentation_date
        date_p.font.size = Pt(14)
        date_p.font.color.rgb = self.theme["accent_light"]
        date_p.alignment = PP_ALIGN.CENTER
        date_p.font.name = self.font_name

        logger.info("[PPT-TITLE] ✅ Title slide created")

    def _agenda_slide(self, agenda_items):
        """
        Create agenda/overview slide
        
        Shows numbered list of topics
        
        Args:
            agenda_items: List of slide titles/topics
        """
        logger.info(f"[PPT-AGENDA] Creating agenda slide with {len(agenda_items)} items")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide, self.theme["primary_light"])

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["primary_dark"]
        header.line.fill.background()

        # Left accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.08), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.theme["accent"]
        accent_bar.line.fill.background()

        # Title
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = "Agenda"
        header_p.font.size = Pt(32)
        header_p.font.bold = True
        header_p.font.color.rgb = self.theme["accent"]
        header_p.font.name = self.font_name

        # Agenda items - numbered list
        agenda_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(8), Inches(5.6))
        agenda_frame = agenda_box.text_frame
        agenda_frame.word_wrap = True

        for i, item in enumerate(agenda_items[:8]):
            if i == 0:
                p = agenda_frame.paragraphs[0]
            else:
                p = agenda_frame.add_paragraph()
            
            p.text = f"{i+1}.  {item}"
            p.level = 0
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = self.theme["primary_dark"]
            p.font.name = self.font_name
            p.space_after = Pt(16)

        logger.info("[PPT-AGENDA] ✅ Agenda slide created")

    def _content_slide(self, title, bullets, image_url=None):
        """
        Create content slide with bullets and optional image
        
        Layout:
        - Header with title
        - Left side: bullet points (5.2 inches if image, 9 inches if no image)
        - Right side: image (if provided)
        
        Args:
            title: Slide title
            bullets: List of bullet points
            image_url: URL to image (optional)
        """
        logger.info(f"[PPT-SLIDE] Creating content slide: '{title}'")
        if image_url:
            logger.info(f"[PPT-SLIDE] Image URL: {str(image_url)[:60]}...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide, self.theme["primary_light"])

        # Header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["primary_dark"]
        header.line.fill.background()

        # Left accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.08), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.theme["accent"]
        accent_bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = self.theme["accent"]
        title_p.font.name = self.font_name

        # Determine layout based on image
        has_image = bool(image_url and str(image_url).strip())
        
        logger.debug(f"[PPT-SLIDE] Has image: {has_image}")

        if has_image:
            # If image: narrow bullets on left
            content_width = Inches(5.2)
            content_left = Inches(0.6)
        else:
            # If no image: full width bullets
            content_width = Inches(9)
            content_left = Inches(0.6)

        # Bullet points
        bullets_box = slide.shapes.add_textbox(content_left, Inches(1.2), content_width, Inches(6))
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        for i, bullet in enumerate(bullets[:6]):  # Max 6 bullets
            if i == 0:
                p = bullets_frame.paragraphs[0]
            else:
                p = bullets_frame.add_paragraph()
            
            # Clean bullet text
            bullet_text = str(bullet).lstrip("•").strip()
            p.text = f"• {bullet_text}"
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme["dark_text"]
            p.font.name = self.font_name
            p.space_after = Pt(14)
            p.line_spacing = 1.15

        # ADD IMAGE TO RIGHT SIDE
        if has_image:
            logger.info(f"[PPT-SLIDE] 🖼️  ATTEMPTING TO ADD IMAGE...")
            
            # Download image
            img_data = self._download_image(image_url)
            
            if img_data:
                try:
                    logger.info(f"[PPT-SLIDE] 📥 Image ready, embedding in slide...")
                    
                    # Add picture to slide
                    picture = slide.shapes.add_picture(
                        img_data,
                        Inches(6.0),      # Left position
                        Inches(1.2),      # Top position
                        width=Inches(3.5),
                        height=Inches(5.8)
                    )
                    
                    logger.info(f"[PPT-SLIDE] ✅ IMAGE EMBEDDED SUCCESSFULLY!")
                    
                except Exception as e:
                    logger.error(f"[PPT-SLIDE] ❌ Error embedding image: {str(e)}", exc_info=True)
            else:
                logger.warning(f"[PPT-SLIDE] ⚠️  Could not download image, showing text-only slide")
        
        logger.info(f"[PPT-SLIDE] ✅ Content slide created: '{title}'")

    def _conclusion_slide(self):
        """
        Create conclusion/key takeaways slide
        
        Shows main insights from presentation
        """
        logger.info("[PPT-CONCLUSION] Creating conclusion slide")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide, self.theme["primary_dark"])

        # Top bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.theme["accent"]
        top_bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = "Key Takeaways"
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = self.theme["accent"]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.name = self.font_name

        # Key points
        points_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(7), Inches(4.5))
        points_frame = points_box.text_frame
        points_frame.word_wrap = True

        default_points = [
            "Comprehensive understanding of core concepts",
            "Practical applications and real-world relevance",
            "Foundation for advanced learning and research",
            "Career opportunities in the field"
        ]

        for i, point in enumerate(default_points):
            if i == 0:
                p = points_frame.paragraphs[0]
            else:
                p = points_frame.add_paragraph()
            
            p.text = f"✓  {point}"
            p.level = 0
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.theme["light_text"]
            p.font.name = self.font_name
            p.space_after = Pt(16)
            p.line_spacing = 1.15

        logger.info("[PPT-CONCLUSION] ✅ Conclusion slide created")

    def _thank_you_slide(self):
        """
        Create closing thank you slide
        
        Professional closing with question prompt
        """
        logger.info("[PPT-THANKYOU] Creating thank you slide")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide, self.theme["accent"])

        # Bar at top
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.15)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["primary_dark"]
        bar.line.fill.background()

        # Main text - "Thank You!"
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = "Thank You!"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.theme["white"]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.font_name

        # Subtitle - "Questions & Discussion"
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        sp = subtitle_frame.paragraphs[0]

        sp.font.size = Pt(28)
        sp.font.color.rgb = self.theme["white"]
        sp.alignment = PP_ALIGN.CENTER
        sp.font.name = self.font_name

        logger.info("[PPT-THANKYOU] ✅ Thank you slide created")

    def generate(self, presentation_data):
        """
        Generate complete presentation from data
        
        Creates:
        1. Title slide
        2. Agenda slide
        3. Content slides (with images)
        4. Conclusion slide
        5. Thank you slide
        
        Args:
            presentation_data: Dictionary containing:
                - title: Presentation title
                - presenter_name: Name of presenter
                - presentation_date: Date
                - subject: Subject/course name
                - slides: List of slide data with title, bullets, image_url
        
        Returns:
            True if successful
        """
        logger.info(f"\n{'='*70}")
        logger.info(f"[PPT-GEN] 🚀 STARTING PRESENTATION GENERATION")
        logger.info(f"[PPT-GEN] {'='*70}\n")
        
        try:
            title = presentation_data.get('title', 'Presentation')
            presenter_name = presentation_data.get('presenter_name', 'Student Name')
            presentation_date = presentation_data.get('presentation_date', 'Academic Year 2025-26')
            subject = presentation_data.get('subject', title)
            slides = presentation_data.get('slides', [])

            logger.info(f"[PPT-GEN] Title: {title}")
            logger.info(f"[PPT-GEN] Presenter: {presenter_name}")
            logger.info(f"[PPT-GEN] Date: {presentation_date}")
            logger.info(f"[PPT-GEN] Total Slides: {len(slides)}\n")

            # SLIDE 1: Title Slide
            logger.info("[PPT-GEN] Creating Slide 1: Title")
            self._title_slide(title, presenter_name, presentation_date, subject)

            # SLIDE 2: Agenda
            if slides:
                logger.info("[PPT-GEN] Creating Slide 2: Agenda")
                agenda_items = [s.get('title', f"Topic {i+1}") for i, s in enumerate(slides[:8])]
                self._agenda_slide(agenda_items)

            # CONTENT SLIDES
            logger.info(f"\n[PPT-GEN] Creating {len(slides)} Content Slides...\n")
            for idx, slide_data in enumerate(slides, 1):
                try:
                    slide_title = slide_data.get('title', f'Slide {idx}')
                    bullets = slide_data.get('bullets', [])
                    image_url = slide_data.get('image_url', None)

                    if not slide_title and not bullets:
                        logger.warning(f"[PPT-GEN] Skipping empty slide {idx}")
                        continue

                    logger.info(f"[PPT-GEN] Creating Content Slide {idx}: {slide_title}")
                    self._content_slide(slide_title, bullets, image_url)
                    
                except Exception as e:
                    logger.error(f"[PPT-GEN] Error on slide {idx}: {str(e)}", exc_info=True)
                    continue

            # CONCLUSION SLIDE
            logger.info("\n[PPT-GEN] Creating Conclusion Slide")
            self._conclusion_slide()

            # THANK YOU SLIDE
            logger.info("[PPT-GEN] Creating Thank You Slide")
            self._thank_you_slide()

            # SAVE FILE
            logger.info(f"\n[PPT-GEN] 💾 Saving presentation...")
            logger.info(f"[PPT-GEN] Output: {self.output_path}")
            self.prs.save(self.output_path)
            
            logger.info(f"\n[PPT-GEN] ✅✅✅ PRESENTATION GENERATED SUCCESSFULLY!")
            logger.info(f"[PPT-GEN] {'='*70}\n")
            
            return True

        except Exception as e:
            logger.error(f"[PPT-GEN] 💥 CRITICAL ERROR: {str(e)}", exc_info=True)
            raise
